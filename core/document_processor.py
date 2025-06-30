"""Document processing functions extracted from main app"""
import os
import tempfile
import hashlib
import textwrap
import time
import logging
import base64
import zipfile
import io
import re
import json as _json
from typing import List, Dict
from openai import AzureOpenAI
import fitz  # PyMuPDF
import pandas as pd
from chunking import DocumentChunker
from tools.aoai import AzureOpenAIClient

def embed_text(oai_client: AzureOpenAI, deployment: str, text: str) -> list[float]:
    """
    Return embedding vector for *text* using the Azure OpenAI client.
    
    This function now uses the proper AzureOpenAIClient from tools.aoai which:
    - Uses actual token counting with tiktoken (not character estimation)
    - Automatically truncates to the 8192 token limit for embedding models
    - Has proper error handling and retry logic
    """
    try:
        # Create an AzureOpenAIClient wrapper for better token handling
        aoai_client = AzureOpenAIClient(document_filename="embedding")
        
        # Use the proper get_embeddings method which handles token limits automatically
        embeddings = aoai_client.get_embeddings(text)
        return embeddings
        
    except Exception as e:
        # Fallback to the existing improved method
        logging.warning(f"AzureOpenAIClient failed, using fallback method: {e}")
        
        MAX_EMBEDDING_TOKENS = 8000  # Conservative limit, leaving 192 tokens buffer
        
        # Try to use tiktoken for accurate token counting if available
        try:
            import tiktoken
            encoding = tiktoken.encoding_for_model("text-embedding-3-large")
            tokens = encoding.encode(text)
            
            if len(tokens) > MAX_EMBEDDING_TOKENS:
                # Truncate at token level for precision
                truncated_tokens = tokens[:MAX_EMBEDDING_TOKENS]
                text = encoding.decode(truncated_tokens)
                logging.warning(f"Text truncated from {len(tokens):,} to {len(truncated_tokens):,} tokens for embedding")
                
        except ImportError:
            # Fallback to character-based estimation if tiktoken not available
            MAX_EMBEDDING_CHARS = 24000  # ~7500 tokens, conservative estimate
            
            if len(text) > MAX_EMBEDDING_CHARS:
                truncated_text = text[:MAX_EMBEDDING_CHARS]
                logging.warning(f"Text truncated from {len(text):,} to {len(truncated_text):,} characters for embedding (tiktoken not available)")
                text = truncated_text
        
        resp = oai_client.embeddings.create(input=[text], model=deployment)
        return resp.data[0].embedding

def pdf_to_documents(pdf_file, oai_client: AzureOpenAI, embed_deployment: str) -> list[dict]:
    """
    Extract pages from an uploaded PDF and return list of documents matching
    the quick‑start index schema.
    """
    docs = []
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(pdf_file.getbuffer())
        tmp_path = tmp.name

    pdf = fitz.open(tmp_path)
    for page_num in range(len(pdf)):
        page_text = pdf[page_num].get_text().strip()
        page_text = f"[{pdf_file.name}] " + page_text
        if not page_text:
            continue
        vector = embed_text(oai_client, embed_deployment, page_text)
        # Build a public-facing link to the PDF (example: blob storage)
        pdf_url = os.getenv("PDF_BASE_URL", "") + pdf_file.name
        doc = {
            "id": hashlib.md5(f"{pdf_file.name}_{page_num}".encode()).hexdigest(),
            "page_chunk": page_text,
            "page_embedding_text_3_large": vector,
            "page_number": page_num + 1,
            "source_file": pdf_file.name,
            "source": pdf_file.name,             # extra alias recognised by agent
            "url": pdf_url,
        }
        docs.append(doc)
    pdf.close()
    return docs

def plainfile_to_docs(
    file_name: str,
    file_bytes: bytes,
    file_url: str,
    oai_client: AzureOpenAI,
    embed_deployment: str,
) -> list[dict]:
    """
    Fallback extractor for DOCX, PPTX, TXT, MD, JSON.
    Tries rich-parsers first; if unavailable, uses a simple XML/text strip.
    Splits the resulting text into ~4 000-char chunks and embeds them.
    """
    ext = os.path.splitext(file_name)[-1].lower()
    txt = ""

    def _strip_xml(xml_bytes: bytes) -> str:
        """Very rough tag-stripper – good enough for plain text."""
        return re.sub(r"<[^>]+>", " ", xml_bytes.decode("utf-8", errors="ignore"))

    try:
        if ext == ".docx":
            try:
                from docx import Document            # python-docx
                txt = "\n".join(p.text for p in Document(io.BytesIO(file_bytes)).paragraphs)
            except ImportError:
                # Fallback: read word/document.xml inside the docx zip
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                    xml = zf.read("word/document.xml")
                txt = _strip_xml(xml)
        elif ext == ".pptx":
            try:
                from pptx import Presentation        # python-pptx
                prs = Presentation(io.BytesIO(file_bytes))
                slides = []
                for slide in prs.slides:
                    slides.append(
                        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                    )
                txt = "\n\n".join(slides)
            except ImportError:
                # Fallback: concatenate text from slide XMLs
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                    slide_files = [n for n in zf.namelist() if n.startswith("ppt/slides/") and n.endswith(".xml")]
                    txt = "\n\n".join(_strip_xml(zf.read(f)) for f in slide_files)
        elif ext in (".txt", ".md"):
            txt = file_bytes.decode("utf-8", errors="ignore")
        elif ext == ".json":
            txt = _json.dumps(_json.loads(file_bytes), indent=2, ensure_ascii=False)
    except Exception as parse_err:
        logging.error("Plain extraction failed for %s: %s", file_name, parse_err)
        txt = ""  # will return [] later if empty

    if not txt.strip():
        return []

    chunks = textwrap.wrap(txt, 4000)
    docs = []
    label = f"[{file_name}] "                 # ← new unified prefix
    for i, chunk in enumerate(chunks):
        chunk_txt = label + chunk             # ← prepend filename
        try:
            vector = embed_text(oai_client, embed_deployment, chunk_txt)   # embed prefixed text
        except Exception as emb_err:
            logging.error("Embedding failed for %s (chunk %d): %s", file_name, i, emb_err)
            continue
        docs.append(
            {
                "id": hashlib.md5(f"{file_name}_{i}".encode()).hexdigest(),
                "page_chunk": chunk_txt,      # store prefixed chunk
                "page_embedding_text_3_large": vector,
                "page_number": i + 1,
                "source_file": file_name,
                "source": file_name,
                "url": file_url or "",
                # Enhanced metadata for fallback processing
                "extraction_method": "simple_parser",
                "document_type": {
                    ".docx": "Word Document",
                    ".pptx": "PowerPoint Presentation", 
                    ".txt": "Text Document",
                    ".md": "Markdown Document",
                    ".json": "JSON Document"
                }.get(ext, "Text Document"),
                "has_figures": False,
                "processing_timestamp": time.strftime('%Y-%m-%dT%H:%M:%SZ', time.gmtime()),
            }
        )
    return docs

def chunk_to_docs(
    file_name: str,
    file_bytes: bytes,
    file_url: str,
    oai_client: AzureOpenAI,
    embed_deployment: str,
) -> list[dict]:
    """
    Run DocumentChunker on *file_bytes* and convert the chunks to the schema
    that agentic-rag indexes (adds url/source_file and ensures embeddings).

    Fix: send **raw bytes** first (needed for XLSX, DOCX, …).  
    Fallback to the old base-64 path if that fails – keeps compatibility
    with formats that still expect a string.
    """
    # Enable multimodal processing for supported file types
    ext = os.path.splitext(file_name)[-1].lower()
    multimodal_env = os.getenv("MULTIMODAL", "false").lower() in ["true", "1", "yes"]
    multimodal_enabled = multimodal_env and ext in ('.pdf', '.png', '.jpeg', '.jpg', '.bmp', '.tiff', '.docx', '.pptx')
    
    dc = DocumentChunker(multimodal=multimodal_enabled, openai_client=oai_client if multimodal_enabled else None)

    def _call_chunker(doc_bytes):
        # DEBUG: Log file size before calling chunker
        if isinstance(doc_bytes, bytes):
            file_size = len(doc_bytes)
            logging.info(f"[chunk_to_docs][{file_name}] Calling chunker with BYTES: {file_size:,} bytes")
            if file_name.endswith('.pdf') and file_size < 1000:
                logging.error(f"[chunk_to_docs][{file_name}] ⚠️ SUSPICIOUS: PDF only {file_size} bytes!")
        elif isinstance(doc_bytes, str):
            logging.info(f"[chunk_to_docs][{file_name}] Calling chunker with STRING: {len(doc_bytes):,} chars")
        else:
            logging.warning(f"[chunk_to_docs][{file_name}] Calling chunker with unknown type: {type(doc_bytes)}")
            
        data = {
            "fileName": file_name,
            "documentBytes": doc_bytes,
            "documentUrl": file_url or "",
        }
        return dc.chunk_documents(data)

    try:
        # DEBUG: Log initial file size
        initial_size = len(file_bytes) if isinstance(file_bytes, bytes) else len(file_bytes) if isinstance(file_bytes, str) else 0
        logging.info(f"[chunk_to_docs][{file_name}] ENTRY POINT: {initial_size:,} bytes/chars, type: {type(file_bytes)}")

        # --- Sanity guard for potentially truncated PDFs ---------------------------------
        if ext == ".pdf":
            header_ok = file_bytes.startswith(b"%PDF-")
            if len(file_bytes) < 2048 or not header_ok:
                logging.warning(
                    f"[chunk_to_docs][{file_name}] 🛑 Suspect truncated/corrupted PDF "
                    f"({len(file_bytes)} bytes, header_ok={header_ok}). "
                    "Falling back to simple pdf_to_documents() extraction."
                )
                raise ValueError("suspect_truncated_pdf")  # triggers fallback section

        chunks, _, _ = _call_chunker(file_bytes)          # 1️⃣ raw bytes
    except Exception as first_err:
        try:
            # 2️⃣ base64 fallback (legacy)
            b64_str = (
                file_bytes
                if isinstance(file_bytes, str)
                else base64.b64encode(file_bytes).decode("utf-8")
            )
            chunks, _, _ = _call_chunker(b64_str)
        except Exception as second_err:
            # 3️⃣ No more fallbacks - all formats handled by DocumentChunker
            # Nothing worked – re-raise original error
            raise first_err from second_err

    docs = []
    label = f"[{file_name}] "                 # ← prefix for DocumentChunker path
    
    # Determine extraction method based on DocumentChunker's chunker selection
    extraction_method = {
        '.vtt': 'transcription_chunker',
        '.json': 'json_chunker', 
        '.xlsx': 'spreadsheet_chunker',
        '.xls': 'spreadsheet_chunker',
        '.pdf': 'document_intelligence',
        '.png': 'document_intelligence',
        '.jpeg': 'document_intelligence', 
        '.jpg': 'document_intelligence',
        '.bmp': 'document_intelligence',
        '.tiff': 'document_intelligence',
        '.docx': 'document_intelligence',
        '.pptx': 'document_intelligence',
        '.nl2sql': 'nl2sql_chunker'
    }.get(ext, 'langchain_chunker')
    
    document_type = {
        '.pdf': 'PDF Document',
        '.docx': 'Word Document', 
        '.pptx': 'PowerPoint Presentation',
        '.xlsx': 'Excel Spreadsheet',
        '.xls': 'Excel Spreadsheet',
        '.csv': 'CSV Spreadsheet',
        '.png': 'Image',
        '.jpg': 'Image',
        '.jpeg': 'Image',
        '.bmp': 'Image',
        '.tiff': 'Image',
        '.html': 'HTML Document',
        '.txt': 'Text Document',
        '.md': 'Markdown Document',
        '.json': 'JSON Document',
        '.vtt': 'Video Transcript',
        '.nl2sql': 'SQL Schema'
    }.get(ext, 'Text Document')
    
    has_figures = False
    
    for i, ch in enumerate(chunks):
        txt = ch.get("page_chunk") or ch.get("chunk") or ch.get("content") or ""
        if not txt:
            continue
        if not txt.startswith(label):         # avoid double-prefix
            txt = label + txt                 # ← prepend filename
            
        # Check if chunk contains figures (for multimodal processing)
        if any(key in ch for key in ['figure_urls', 'figure_descriptions', 'combined_caption', 'relatedImages', 'isMultimodal']):
            has_figures = True
            
        # embedding – reuse if present, else create with safe fallback
        vector = ch.get("page_embedding_text_3_large")
        if not vector:
            try:
                vector = embed_text(oai_client, embed_deployment, txt)
            except Exception as emb_err:
                logging.error("Embedding failed for %s (chunk %d): %s", file_name, i, emb_err)
                continue  # skip this chunk
        
        # Safe join for image captions - handle case where captions might be dicts instead of strings
        def _safe_join_captions(captions_list):
            safe_captions = []
            if captions_list:
                for caption in captions_list:
                    if isinstance(caption, str):
                        safe_captions.append(caption)
                    elif isinstance(caption, dict):
                        # If it's a dict, try to extract meaningful content
                        if "content" in caption:
                            safe_captions.append(str(caption["content"]))
                        else:
                            safe_captions.append(str(caption))
                    else:
                        safe_captions.append(str(caption))
            return " ".join(safe_captions) if safe_captions else ""

        # Generate caption vector for multimodal content
        caption_vector = [0.0] * 3072  # Default to zero vector to avoid null value errors (text-embedding-3-large dims)
        image_captions = ch.get("imageCaptions", [])
        if image_captions and isinstance(image_captions, list) and len(image_captions) > 0:
            # Use safe join to handle mixed types (strings, dicts, etc.)
            captions_text = _safe_join_captions(image_captions)
            if captions_text:
                try:
                    caption_vector = embed_text(oai_client, embed_deployment, captions_text)
                except Exception as emb_err:
                    logging.error("Caption embedding failed for %s (chunk %d): %s", file_name, i, emb_err)
                    # Keep the zero vector as fallback
        
        docs.append(
            {
                "id": ch.get("id") or hashlib.md5(f"{file_name}_{i}".encode()).hexdigest(),
                "page_chunk": txt,
                "page_embedding_text_3_large": vector,
                "content": ch.get("content", txt),  # Use multimodal content if available, fallback to txt
                "contentVector": ch.get("contentVector") or vector,  # Use multimodal vector if available
                "page_number": ch.get("page_number") or i + 1,
                "source_file": file_name,
                "source": file_name,
                "url": file_url or "",
                # Enhanced metadata
                "extraction_method": extraction_method,
                "document_type": document_type, 
                "has_figures": has_figures,
                "processing_timestamp": time.strftime('%Y-%m-%dT%H:%M:%SZ', time.gmtime()),
                # Multimodal fields (if present in chunk)
                "imageCaptions": _safe_join_captions(image_captions),
                "captionVector": caption_vector,
                "relatedImages": ch.get("relatedImages", []),
                "isMultimodal": ch.get("isMultimodal", bool(image_captions or ch.get("relatedImages"))),
                "filename": ch.get("filename", file_name),
            }
        )
    return docs

def tabular_to_docs(
    file_name: str,
    file_bytes: bytes,
    file_url: str,
    oai_client: AzureOpenAI,
    embed_deployment: str,
) -> list[dict]:
    """
    Last-resort converter for CSV/XLS/XLSX when DocumentChunker fails.
    Turns every ~4 000-char slice of the table (as CSV text) into one doc.
    """
    ext = os.path.splitext(file_name)[-1].lower()
    # Extract plain text
    if file_name.lower().endswith(".csv"):
        txt = file_bytes.decode("utf-8", errors="ignore")
    else:  # .xls / .xlsx
        excel = pd.ExcelFile(io.BytesIO(file_bytes))
        txt_parts = []
        for sheet in excel.sheet_names:
            df = excel.parse(sheet)
            txt_parts.append(f"\n\n### Sheet: {sheet}\n" + df.to_csv(index=False))
        txt = "".join(txt_parts)

    chunks = textwrap.wrap(txt, 4000)
    docs = []
    label = f"[{file_name}] "                 # ← prefix for CSV / Excel
    for i, chunk in enumerate(chunks):
        chunk_txt = label + chunk
        try:
            vector = embed_text(oai_client, embed_deployment, chunk_txt)
        except Exception as emb_err:
            logging.error("Embedding failed for %s (chunk %d): %s", file_name, i, emb_err)
            continue  # skip this chunk
        docs.append(
            {
                "id": hashlib.md5(f"{file_name}_{i}".encode()).hexdigest(),
                "page_chunk": chunk_txt,
                "page_embedding_text_3_large": vector,
                "page_number": i + 1,
                "source_file": file_name,
                "source": file_name,
                "url": file_url or "",
                # Enhanced metadata for tabular processing
                "extraction_method": "pandas_parser",
                "document_type": {
                    ".csv": "CSV Spreadsheet",
                    ".xlsx": "Excel Spreadsheet", 
                    ".xls": "Excel Spreadsheet"
                }.get(ext, "Tabular Data"),
                "has_figures": False,
                "processing_timestamp": time.strftime('%Y-%m-%dT%H:%M:%SZ', time.gmtime()),
            }
        )
    return docs
