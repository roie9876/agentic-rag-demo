#!/usr/bin/env python3
"""
agentic-rag-demo.py
===================
CLI demo of Agentic Retrieval‑Augmented Generation on Azure
compatible with **openai‑python ≥ 1.0**.

Based on the official quick‑start:
https://learn.microsoft.com/azure/search/search-get-started-agentic-retrieval?pivots=python
"""
from __future__ import annotations

import argparse
import json
import os
import sys
import textwrap
import subprocess   # for az cli calls
import httpx  # HTTP probe for RBAC status
import zipfile
import tempfile
import subprocess
import time  # Added to support sleep in polling after ingestion
import base64  # Added to encode document bytes for _chunk_to_docs

from pathlib import Path
from typing import List, Tuple, Dict

import pandas as pd           # ← ADD THIS LINE
# ---------------------------------------------------------------------------
# Streamlit Data‑Editor helper (works on both old & new versions)
# ---------------------------------------------------------------------------
import streamlit as st
from chunking import DocumentChunker
from tools.aoai import AzureOpenAIClient

# Import the test_retrieval module
from test_retrieval import render_test_retrieval_tab

def _st_data_editor(*args, **kwargs):
    """
    Wrapper that tries st.data_editor (Streamlit ≥ 1.29) and falls back to
    st.experimental_data_editor for older releases.
    """
    if hasattr(st, "data_editor"):
        return st.data_editor(*args, **kwargs)
    elif hasattr(st, "experimental_data_editor"):
        return st.experimental_data_editor(*args, **kwargs)
    else:
        st.error(
            "⚠️ Your Streamlit version is too old for data‑editor. "
            "Upgrade with:\n\n"
            "    pip install --upgrade streamlit"
        )
        st.stop()
# Reliable check whether code runs under `streamlit run …`
try:
    from streamlit.runtime import exists as _st_in_runtime
except ImportError:       # fallback for older Streamlit
    _st_in_runtime = lambda: False

import re  # for citation parsing
from azure.identity import DefaultAzureCredential
from openai import AzureOpenAI
from azure.core.credentials import AzureKeyCredential
from azure.core.exceptions import ClientAuthenticationError, HttpResponseError
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient  # NEW
from azure.search.documents.indexes.models import (
    SearchIndex,
    SimpleField,
    SearchFieldDataType,
    SearchableField,
    SearchField,
    VectorSearch,
    HnswAlgorithmConfiguration,
    VectorSearchProfile,
    SemanticConfiguration,
    SemanticField,
    SemanticPrioritizedFields,
    SemanticSearch,
    AzureOpenAIVectorizer,
    AzureOpenAIVectorizerParameters,
    KnowledgeAgent,
    KnowledgeAgentAzureOpenAIModel,
    KnowledgeAgentTargetIndex,
    KnowledgeAgentRequestLimits,
)

# --- Logging setup ---
import logging
import inspect
# Show only warnings and errors in the terminal
logging.basicConfig(level=logging.WARNING, format="%(levelname)s: %(message)s")

# Knowledge‑agent runtime
from azure.search.documents.agent import KnowledgeAgentRetrievalClient
from azure.search.documents.agent.models import (
    KnowledgeAgentRetrievalRequest,
    KnowledgeAgentMessage,
    KnowledgeAgentMessageTextContent,
    KnowledgeAgentIndexParams,
)
from dotenv import load_dotenv

# Health Check Module
from health_check import HealthChecker, HealthCheckUI

# Import the new module for AI Foundry functionality
from agent_foundry import (
    check_azure_cli_login, 
    get_ai_foundry_projects,
    create_ai_foundry_agent
)

# Import Azure Function helper module
from azure_function_helper import (
    load_env_vars,
    get_azure_subscription,
    list_function_apps,
    load_function_settings,
    push_function_settings,
    deploy_function_code
)

from azure.search.documents import SearchIndexingBufferedSender  # NEW
import fitz                            # PyMuPDF
import hashlib, tempfile               # for PDF processing
import zipfile
from azure.identity import AzureCliCredential, get_bearer_token_provider
from azure.ai.projects import AIProjectClient
# --- Azure AI Foundry SDK ----------------------------------------------------
# Support multiple SDK generations where the tool classes moved packages/names
try:
    # GA / recent preview: everything under azure.ai.agents
    from azure.ai.agents import FunctionTool, FunctionDefinition
except ImportError:
    # Older builds may expose FunctionTool and/or FunctionDefinition
    # under azure.ai.agents.models
    try:
        from azure.ai.agents import FunctionTool  # type: ignore
    except ImportError:
        FunctionTool = None  # type: ignore
    try:
        from azure.ai.agents.models import FunctionTool as _FTModel, FunctionDefinition  # type: ignore
        if FunctionTool is None:  # fallback when only the models version exists
            FunctionTool = _FTModel  # type: ignore
    except ImportError:
        FunctionDefinition = None  # type: ignore

# OpenAPI tool helper (available in azure‑ai‑agents ≥ 1.0.0b2)
from azure.ai.agents.models import OpenApiTool, OpenApiAnonymousAuthDetails


# ToolDefinition is only under .models
def _search_credential() -> AzureKeyCredential | DefaultAzureCredential:
    """
    Return Azure credential based on env:
    • If AZURE_SEARCH_KEY is set → key auth
    • else → DefaultAzureCredential (AAD)
    """
    key = os.getenv("AZURE_SEARCH_KEY", "").strip()
    if key:
        return AzureKeyCredential(key)
    return DefaultAzureCredential()


# ---------------------------------------------------------------------------
# RBAC status probe
# ---------------------------------------------------------------------------
def _rbac_enabled(service_url: str) -> bool:
    """
    Quick probe: return True if Role‑based access control is enabled on the
    Search service (Authentication mode = RBAC).
    """
    try:
        url = f"{service_url.rstrip('/')}/?api-version=2023-11-01"
        r = httpx.get(url, timeout=3)
        # When RBAC is ON the payload includes "RoleBasedAccessControl"
        return "RoleBasedAccessControl" in r.text
    except Exception:
        return False

# ---------------------------------------------------------------------------
# Azure CLI helpers
# ---------------------------------------------------------------------------
def _az_logged_user() -> tuple[str | None, str | None]:
    """Return (UPN/email, subscription-id) of the signed‑in az cli user, or (None,None)."""
    try:
        out = subprocess.check_output(
            ["az", "account", "show", "--output", "json"], text=True, timeout=3
        )
        data = json.loads(out)
        return data["user"]["name"], data["id"]
    except Exception:
        return None, None


# Remove the duplicate function since it's now in agent_foundry.py
# def check_azure_cli_login() -> tuple[bool, dict | None]:
#     """
#     Return (logged_in, account_json_or_None) by running `az account show`.
#     """
#     try:
#         out = subprocess.check_output(
#             ["az", "account", "show", "--output", "json"],
#             text=True,
#             timeout=5,
#         )
#         return True, json.loads(out)
#     except subprocess.CalledProcessError:
#         return False, None
#     except Exception:
#         return False, None


# Remove the duplicate function since it's now in agent_foundry.py
# def get_ai_foundry_projects(cred: AzureCliCredential) -> list[dict]:
#     """
#     Return a list of Foundry projects visible to the signed‑in CLI user via
#     `az ai project list`. Each item includes:
#         {name, location, endpoint, resource_group, hub_name}
#     """
#     try:
#         out = subprocess.check_output(
#             ["az", "ai", "project", "list", "--output", "json"],
#             text=True,
#             timeout=10,
#         )
#         data = json.loads(out)
#         projs = []
#         for p in data:
#             projs.append(
#                 {
#                     "name": p["name"],
#                     "location": p["location"],
#                     "endpoint": p["properties"]["endpoint"],
#                     "resource_group": p["resourceGroup"],
#                     "hub_name": p["properties"].get("hubName", ""),
#                 }
#             )
#         return projs
#     except Exception as err:
#         logging.warning("Failed to list AI Foundry projects: %s", err)
#         return []


def _grant_search_role(service_name: str, subscription_id: str, resource_group: str, principal: str, role: str) -> tuple[bool, str]:
    """
    Grant the specified *role* to *principal* on the given service.
    Returns (success, message).
    """
    try:
        scope = f"/subscriptions/{subscription_id}/resourceGroups/{resource_group}/providers/Microsoft.Search/searchServices/{service_name}"
        subprocess.check_call(
            [
                "az", "role", "assignment", "create",
                "--role", role,
                "--assignee", principal,
                "--scope", scope
            ],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=15,
        )
        return True, "Role granted successfully"
    except subprocess.CalledProcessError as e:
        return False, f"az cli error: {e}"
    except Exception as ex:
        return False, str(ex)



# ---------------------------------------------------------------------------
# Helper to grant OpenAI role
# ---------------------------------------------------------------------------
def _grant_openai_role(account_name: str, subscription_id: str, resource_group: str,
                       principal: str, role: str = "Cognitive Services OpenAI User") -> tuple[bool, str]:
    """
    Grant *role* (default: Cognitive Services OpenAI User) on the Azure OpenAI
    account to *principal*.
    """
    try:
        scope = (
            f"/subscriptions/{subscription_id}/resourceGroups/{resource_group}"
            f"/providers/Microsoft.CognitiveServices/accounts/{account_name}"
        )
        subprocess.check_call(
            [
                "az", "role", "assignment", "create",
                "--role", role,
                "--assignee", principal,
                "--scope", scope,
            ],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=15,
        )
        return True, "Role granted successfully"
    except subprocess.CalledProcessError as e:
        return False, f"az cli error: {e}"
    except Exception as ex:
        return False, str(ex)

# ---------------------------------------------------------------------------
# Force‑reload .env at runtime
# ---------------------------------------------------------------------------
def _reload_env_and_restart():
    """
    Reload the .env file (override existing variables), clear cached clients,
    and rerun the Streamlit script so the new values take effect.
    """
    env_path = Path(__file__).resolve().parent / ".env"
    load_dotenv(env_path, override=True)

    # Clear Streamlit caches for resource‑building functions
    for fn in (init_openai, init_search_client, init_agent_client):
        if hasattr(fn, "clear"):
            fn.clear()

    st.toast("✅ .env reloaded – restarting app…", icon="🔄")
    if hasattr(st, "rerun"):
        st.rerun()
    else:  # fallback for older versions
        st.experimental_rerun()

##############################################################################
# Environment helpers
##############################################################################

load_dotenv(Path(__file__).resolve().parent / ".env")
FUNCTION_KEY = os.getenv("AGENT_FUNC_KEY", "")
TOP_K_DEFAULT = int(os.getenv("TOP_K", 5))   # fallback for CLI path


def env(var: str) -> str:
    """Fetch env var or exit with error."""
    v = os.getenv(var)
    if not v:
        sys.exit(f"❌ Missing env var: {var}")
    return v


def init_openai(model: str = "o3") -> Tuple[AzureOpenAI, dict]:
    """
    Return AzureOpenAI client + chat params for the chosen *model*
    (o3, 4o, 41).  Picks env‑vars with appropriate suffix.
    """
    suffix = {"o3": "", "4o": "_4o", "41": "_41"}.get(model, "")
    client = AzureOpenAI(
        api_key=env(f"AZURE_OPENAI_KEY{suffix}"),
        azure_endpoint=env(f"AZURE_OPENAI_ENDPOINT{suffix}").rstrip("/"),
        api_version=env(f"AZURE_OPENAI_API_VERSION{suffix}"),
    )
    # If no API key, use AAD token provider
    if not os.getenv(f"AZURE_OPENAI_KEY{suffix}", "").strip():
        # Switch to AAD token provider
        aad = get_bearer_token_provider(
            DefaultAzureCredential(), "https://cognitiveservices.azure.com/.default"
        )
        client = AzureOpenAI(
            azure_endpoint=env(f"AZURE_OPENAI_ENDPOINT{suffix}").rstrip("/"),
            azure_ad_token_provider=aad,
            api_version=env(f"AZURE_OPENAI_API_VERSION{suffix}"),
        )
    chat_params = dict(
        model=env(f"AZURE_OPENAI_DEPLOYMENT{suffix}"),
        temperature=0,
        max_tokens=st.session_state.get("max_tokens", 80000),
    )
    return client, chat_params


def embed_text(oai_client: AzureOpenAI, deployment: str, text: str) -> list[float]:
    """
    Return embedding vector for *text* using the Azure OpenAI client.
    
    This function now uses the proper AzureOpenAIClient from tools.aoai which:
    - Uses actual token counting with tiktoken (not character estimation)
    - Automatically truncates to the 8192 token limit for embedding models
    - Has proper error handling and retry logic
    """
    try:
        # Create an AzureOpenAIClient wrapper for better token handling
        aoai_client = AzureOpenAIClient(document_filename="embedding")
        
        # Use the proper get_embeddings method which handles token limits automatically
        embeddings = aoai_client.get_embeddings(text)
        return embeddings
        
    except Exception as e:
        # Fallback to the existing improved method
        logging.warning(f"AzureOpenAIClient failed, using fallback method: {e}")
        
        MAX_EMBEDDING_TOKENS = 8000  # Conservative limit, leaving 192 tokens buffer
        
        # Try to use tiktoken for accurate token counting if available
        try:
            import tiktoken
            encoding = tiktoken.encoding_for_model("text-embedding-3-large")
            tokens = encoding.encode(text)
            
            if len(tokens) > MAX_EMBEDDING_TOKENS:
                # Truncate at token level for precision
                truncated_tokens = tokens[:MAX_EMBEDDING_TOKENS]
                text = encoding.decode(truncated_tokens)
                logging.warning(f"Text truncated from {len(tokens):,} to {len(truncated_tokens):,} tokens for embedding")
                
        except ImportError:
            # Fallback to character-based estimation if tiktoken not available
            MAX_EMBEDDING_CHARS = 24000  # ~7500 tokens, conservative estimate
            
            if len(text) > MAX_EMBEDDING_CHARS:
                truncated_text = text[:MAX_EMBEDDING_CHARS]
                logging.warning(f"Text truncated from {len(text):,} to {len(truncated_text):,} characters for embedding (tiktoken not available)")
                text = truncated_text
        
        resp = oai_client.embeddings.create(input=[text], model=deployment)
        return resp.data[0].embedding


def pdf_to_documents(pdf_file, oai_client: AzureOpenAI, embed_deployment: str) -> list[dict]:
    """
    Extract pages from an uploaded PDF and return list of documents matching
    the quick‑start index schema.
    """
    docs = []
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(pdf_file.getbuffer())
        tmp_path = tmp.name

    pdf = fitz.open(tmp_path)
    for page_num in range(len(pdf)):
        page_text = pdf[page_num].get_text().strip()
        page_text = f"[{pdf_file.name}] " + page_text
        if not page_text:
            continue
        vector = embed_text(oai_client, embed_deployment, page_text)
        # Build a public-facing link to the PDF (example: blob storage)
        pdf_url = os.getenv("PDF_BASE_URL", "") + pdf_file.name
        doc = {
            "id": hashlib.md5(f"{pdf_file.name}_{page_num}".encode()).hexdigest(),
            "page_chunk": page_text,
            "page_embedding_text_3_large": vector,
            "page_number": page_num + 1,
            "source_file": pdf_file.name,
            "source": pdf_file.name,             # extra alias recognised by agent
            "url": pdf_url,
        }
        docs.append(doc)
    pdf.close()
    return docs


# ───────── Text-based fallback for Office / plain files ─────────
def _plainfile_to_docs(
    file_name: str,
    file_bytes: bytes,
    file_url: str,
    oai_client: AzureOpenAI,
    embed_deployment: str,
) -> list[dict]:
    """
    Fallback extractor for DOCX, PPTX, TXT, MD, JSON.
    Tries rich-parsers first; if unavailable, uses a simple XML/text strip.
    Splits the resulting text into ~4 000-char chunks and embeds them.
    """
    import re, zipfile, io
    ext = os.path.splitext(file_name)[-1].lower()
    txt = ""

    def _strip_xml(xml_bytes: bytes) -> str:
        """Very rough tag-stripper – good enough for plain text."""
        return re.sub(r"<[^>]+>", " ", xml_bytes.decode("utf-8", errors="ignore"))

    try:
        if ext == ".docx":
            try:
                from docx import Document            # python-docx
                txt = "\n".join(p.text for p in Document(io.BytesIO(file_bytes)).paragraphs)
            except ImportError:
                # Fallback: read word/document.xml inside the docx zip
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                    xml = zf.read("word/document.xml")
                txt = _strip_xml(xml)
        elif ext == ".pptx":
            try:
                from pptx import Presentation        # python-pptx
                prs = Presentation(io.BytesIO(file_bytes))
                slides = []
                for slide in prs.slides:
                    slides.append(
                        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                    )
                txt = "\n\n".join(slides)
            except ImportError:
                # Fallback: concatenate text from slide XMLs
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                    slide_files = [n for n in zf.namelist() if n.startswith("ppt/slides/") and n.endswith(".xml")]
                    txt = "\n\n".join(_strip_xml(zf.read(f)) for f in slide_files)
        elif ext in (".txt", ".md"):
            txt = file_bytes.decode("utf-8", errors="ignore")
        elif ext == ".json":
            import json as _json
            txt = _json.dumps(_json.loads(file_bytes), indent=2, ensure_ascii=False)
    except Exception as parse_err:
        logging.error("Plain extraction failed for %s: %s", file_name, parse_err)
        txt = ""  # will return [] later if empty

    if not txt.strip():
        return []

    chunks = textwrap.wrap(txt, 4000)
    docs = []
    label = f"[{file_name}] "                 # ← new unified prefix
    for i, chunk in enumerate(chunks):
        chunk_txt = label + chunk             # ← prepend filename
        try:
            vector = embed_text(oai_client, embed_deployment, chunk_txt)   # embed prefixed text
        except Exception as emb_err:
            logging.error("Embedding failed for %s (chunk %d): %s", file_name, i, emb_err)
            continue
        docs.append(
            {
                "id": hashlib.md5(f"{file_name}_{i}".encode()).hexdigest(),
                "page_chunk": chunk_txt,      # store prefixed chunk
                "page_embedding_text_3_large": vector,
                "page_number": i + 1,
                "source_file": file_name,
                "source": file_name,
                "url": file_url or "",
                # Enhanced metadata for fallback processing
                "extraction_method": "simple_parser",
                "document_type": {
                    ".docx": "Word Document",
                    ".pptx": "PowerPoint Presentation", 
                    ".txt": "Text Document",
                    ".md": "Markdown Document",
                    ".json": "JSON Document"
                }.get(ext, "Text Document"),
                "has_figures": False,
                "processing_timestamp": time.strftime('%Y-%m-%dT%H:%M:%SZ', time.gmtime()),
            }
        )
    return docs
# -----------------------------------------------------------------

def _chunk_to_docs(
    file_name: str,
    file_bytes: bytes,
    file_url: str,
    oai_client: AzureOpenAI,
    embed_deployment: str,
) -> list[dict]:
    """
    Run DocumentChunker on *file_bytes* and convert the chunks to the schema
    that agentic-rag indexes (adds url/source_file and ensures embeddings).

    Fix: send **raw bytes** first (needed for XLSX, DOCX, …).  
    Fallback to the old base-64 path if that fails – keeps compatibility
    with formats that still expect a string.
    """
    # Enable multimodal processing for supported file types
    ext = os.path.splitext(file_name)[-1].lower()
    multimodal_env = os.getenv("MULTIMODAL", "false").lower() in ["true", "1", "yes"]
    multimodal_enabled = multimodal_env and ext in ('.pdf', '.png', '.jpeg', '.jpg', '.bmp', '.tiff', '.docx', '.pptx')
    
    dc = DocumentChunker(multimodal=multimodal_enabled, openai_client=oai_client if multimodal_enabled else None)

    ext = os.path.splitext(file_name)[-1].lower()
    # ── EARLY BYPASS for known troublesome formats ────────────────
    if ext in (".csv", ".xls", ".xlsx"):
        return _tabular_to_docs(file_name, file_bytes, file_url, oai_client, embed_deployment)
    # Remove DOCX/PPTX from bypass to allow Document Intelligence processing
    if ext in (".txt", ".md", ".json"):
        return _plainfile_to_docs(file_name, file_bytes, file_url, oai_client, embed_deployment)
    # ------------------------------------------------------------------

    def _call_chunker(doc_bytes):
        # DEBUG: Log file size before calling chunker
        if isinstance(doc_bytes, bytes):
            file_size = len(doc_bytes)
            logging.info(f"[_chunk_to_docs][{file_name}] Calling chunker with BYTES: {file_size:,} bytes")
            if file_name.endswith('.pdf') and file_size < 1000:
                logging.error(f"[_chunk_to_docs][{file_name}] ⚠️ SUSPICIOUS: PDF only {file_size} bytes!")
        elif isinstance(doc_bytes, str):
            logging.info(f"[_chunk_to_docs][{file_name}] Calling chunker with STRING: {len(doc_bytes):,} chars")
        else:
            logging.warning(f"[_chunk_to_docs][{file_name}] Calling chunker with unknown type: {type(doc_bytes)}")
            
        data = {
            "fileName": file_name,
            "documentBytes": doc_bytes,
            "documentUrl": file_url or "",
        }
        return dc.chunk_documents(data)

    try:
        # DEBUG: Log initial file size
        initial_size = len(file_bytes) if isinstance(file_bytes, bytes) else len(file_bytes) if isinstance(file_bytes, str) else 0
        logging.info(f"[_chunk_to_docs][{file_name}] ENTRY POINT: {initial_size:,} bytes/chars, type: {type(file_bytes)}")

        # --- Sanity guard for potentially truncated PDFs ---------------------------------
        if ext == ".pdf":
            header_ok = file_bytes.startswith(b"%PDF-")
            if len(file_bytes) < 2048 or not header_ok:
                logging.warning(
                    f"[_chunk_to_docs][{file_name}] 🛑 Suspect truncated/corrupted PDF "
                    f"({len(file_bytes)} bytes, header_ok={header_ok}). "
                    "Falling back to simple pdf_to_documents() extraction."
                )
                raise ValueError("suspect_truncated_pdf")  # triggers fallback section

        chunks, _, _ = _call_chunker(file_bytes)          # 1️⃣ raw bytes
    except Exception as first_err:
        try:
            # 2️⃣ base64 fallback (legacy)
            b64_str = (
                file_bytes
                if isinstance(file_bytes, str)
                else base64.b64encode(file_bytes).decode("utf-8")
            )
            chunks, _, _ = _call_chunker(b64_str)
        except Exception as second_err:
            # 3️⃣ tabular fallback for XLS/CSV
            ext = os.path.splitext(file_name)[-1].lower()
            if ext in (".csv", ".xlsx", ".xls"):
                return _tabular_to_docs(file_name, file_bytes, file_url, oai_client, embed_deployment)
            # Nothing worked – re-raise original error
            raise first_err from second_err

    docs = []
    label = f"[{file_name}] "                 # ← prefix for DocumentChunker path
    
    # Determine extraction method and document type
    extraction_method = "document_intelligence" if ext in ('.pdf', '.png', '.jpeg', '.jpg', '.bmp', '.tiff', '.docx', '.pptx', '.xlsx', '.html') else "langchain_chunker"
    document_type = {
        '.pdf': 'PDF Document',
        '.docx': 'Word Document', 
        '.pptx': 'PowerPoint Presentation',
        '.xlsx': 'Excel Spreadsheet',
        '.png': 'Image',
        '.jpg': 'Image',
        '.jpeg': 'Image',
        '.bmp': 'Image',
        '.tiff': 'Image',
        '.html': 'HTML Document'
    }.get(ext, 'Text Document')
    
    has_figures = False
    
    # Check if any chunk indicates a corrupted file - REMOVED since we no longer validate PDF headers
    # Let Document Intelligence handle all format validation
    
    
    for i, ch in enumerate(chunks):
        txt = ch.get("page_chunk") or ch.get("chunk") or ch.get("content") or ""
        if not txt:
            continue
        if not txt.startswith(label):         # avoid double-prefix
            txt = label + txt                 # ← prepend filename
            
        # Check if chunk contains figures (for multimodal processing)
        if any(key in ch for key in ['figure_urls', 'figure_descriptions', 'combined_caption', 'relatedImages', 'isMultimodal']):
            has_figures = True
            
        # embedding – reuse if present, else create with safe fallback
        vector = ch.get("page_embedding_text_3_large")
        if not vector:
            try:
                vector = embed_text(oai_client, embed_deployment, txt)
            except Exception as emb_err:
                logging.error("Embedding failed for %s (chunk %d): %s", file_name, i, emb_err)
                continue  # skip this chunk
        
        # Safe join for image captions - handle case where captions might be dicts instead of strings
        def _safe_join_captions(captions_list):
            safe_captions = []
            if captions_list:
                for caption in captions_list:
                    if isinstance(caption, str):
                        safe_captions.append(caption)
                    elif isinstance(caption, dict):
                        # If it's a dict, try to extract meaningful content
                        if "content" in caption:
                            safe_captions.append(str(caption["content"]))
                        else:
                            safe_captions.append(str(caption))
                    else:
                        safe_captions.append(str(caption))
            return " ".join(safe_captions) if safe_captions else ""

        # Generate caption vector for multimodal content
        caption_vector = None
        image_captions = ch.get("imageCaptions", [])
        if image_captions and isinstance(image_captions, list) and len(image_captions) > 0:
            # Use safe join to handle mixed types (strings, dicts, etc.)
            captions_text = _safe_join_captions(image_captions)
            if captions_text:
                try:
                    caption_vector = embed_text(oai_client, embed_deployment, captions_text)
                except Exception as emb_err:
                    logging.error("Caption embedding failed for %s (chunk %d): %s", file_name, i, emb_err)
        
        docs.append(
            {
                "id": ch.get("id") or hashlib.md5(f"{file_name}_{i}".encode()).hexdigest(),
                "page_chunk": txt,
                "page_embedding_text_3_large": vector,
                "content": ch.get("content", txt),  # Use multimodal content if available, fallback to txt
                "contentVector": ch.get("contentVector") or vector,  # Use multimodal vector if available
                "page_number": ch.get("page_number") or i + 1,
                "source_file": file_name,
                "source": file_name,
                "url": file_url or "",
                # Enhanced metadata
                "extraction_method": extraction_method,
                "document_type": document_type, 
                "has_figures": has_figures,
                "processing_timestamp": time.strftime('%Y-%m-%dT%H:%M:%SZ', time.gmtime()),
                # Multimodal fields (if present in chunk)
                "imageCaptions": _safe_join_captions(image_captions),
                "captionVector": caption_vector,
                "relatedImages": ch.get("relatedImages", []),
                "isMultimodal": ch.get("isMultimodal", bool(image_captions or ch.get("relatedImages"))),
                "filename": ch.get("filename", file_name),
            }
        )
    return docs
# ───────────────────────── end helper ────────────────────────────────


# ─────────────── Fallback: CSV / XLS(X) → docs ───────────────
def _tabular_to_docs(
    file_name: str,
    file_bytes: bytes,
    file_url: str,
    oai_client: AzureOpenAI,
    embed_deployment: str,
) -> list[dict]:
    """
    Last-resort converter for CSV/XLS/XLSX when DocumentChunker fails.
    Turns every ~4 000-char slice of the table (as CSV text) into one doc.
    """
    import io, pandas as pd  # pandas is only needed here
    ext = os.path.splitext(file_name)[-1].lower()  # Add ext definition
    # Extract plain text
    if file_name.lower().endswith(".csv"):
        txt = file_bytes.decode("utf-8", errors="ignore")
    else:  # .xls / .xlsx
        excel = pd.ExcelFile(io.BytesIO(file_bytes))
        txt_parts = []
        for sheet in excel.sheet_names:
            df = excel.parse(sheet)
            txt_parts.append(f"\n\n### Sheet: {sheet}\n" + df.to_csv(index=False))
        txt = "".join(txt_parts)

    chunks = textwrap.wrap(txt, 4000)
    docs = []
    label = f"[{file_name}] "                 # ← prefix for CSV / Excel
    for i, chunk in enumerate(chunks):
        chunk_txt = label + chunk
        try:
            vector = embed_text(oai_client, embed_deployment, chunk_txt)
        except Exception as emb_err:
            logging.error("Embedding failed for %s (chunk %d): %s", file_name, i, emb_err)
            continue  # skip this chunk
        docs.append(
            {
                "id": hashlib.md5(f"{file_name}_{i}".encode()).hexdigest(),
                "page_chunk": chunk_txt,
                "page_embedding_text_3_large": vector,
                "page_number": i + 1,
                "source_file": file_name,
                "source": file_name,
                "url": file_url or "",
                # Enhanced metadata for tabular processing
                "extraction_method": "pandas_parser",
                "document_type": {
                    ".csv": "CSV Spreadsheet",
                    ".xlsx": "Excel Spreadsheet", 
                    ".xls": "Excel Spreadsheet"
                }.get(ext, "Tabular Data"),
                "has_figures": False,
                "processing_timestamp": time.strftime('%Y-%m-%dT%H:%M:%SZ', time.gmtime()),
            }
        )
    return docs


def init_search_client(index_name: str | None = None) -> Tuple[SearchClient, SearchIndexClient]:
    """
    Return (search_client, index_client).
    Only API Key authentication is supported for agentic retrieval (see Azure docs).
    `index_name` – if provided, SearchClient will target that index,
    otherwise a dummy client pointing at the service root is returned.
    """
    endpoint = env("AZURE_SEARCH_ENDPOINT")
    credential = _search_credential()

    index_client = SearchIndexClient(endpoint=endpoint, credential=credential)

    if index_name:
        search_client = SearchClient(endpoint=endpoint, index_name=index_name, credential=credential)
    else:
        # Return a SearchClient bound to some default index just to keep type happy.
        # Caller should create a proper one later.
        search_client = SearchClient(endpoint=endpoint, index_name="dummy", credential=credential)

    # Improved index listing debug
    try:
        available_indexes = list(index_client.list_indexes())
        st.session_state.available_indexes = [idx.name for idx in available_indexes]
    except Exception as conn_error:
        logging.error("list_indexes() failed: %s", conn_error)
        st.session_state.available_indexes = []

    return search_client, index_client

# ---------------------------------------------------------------------------
# Knowledge‑Agent client (cached per agent name)
# ---------------------------------------------------------------------------
@st.cache_resource
def init_agent_client(agent_name: str) -> KnowledgeAgentRetrievalClient:
    """
    Create KnowledgeAgentRetrievalClient.
    Only API Key authentication is supported for agentic retrieval (see Azure docs).
    """
    cred = _search_credential()   # Always AzureKeyCredential
    return KnowledgeAgentRetrievalClient(
        endpoint=env("AZURE_SEARCH_ENDPOINT"),
        agent_name=agent_name,
        credential=cred,
    )


def create_agentic_rag_index(index_client: "SearchIndexClient", name: str) -> bool:
    """
    Create (or recreate) an index + knowledge-agent עם מפתח-API ל-Azure OpenAI.
    """
    try:
        # ----------- הגדרות בסיסיות -----------------
        azure_openai_endpoint = env("AZURE_OPENAI_ENDPOINT_41")
        embedding_deployment  = os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT", "text-embedding-3-large")
        embedding_model       = os.getenv("AZURE_OPENAI_EMBEDDING_MODEL",      "text-embedding-3-large")
        VECTOR_DIM            = 3072
        # Resolve OpenAI key – prefer suffix _41, fall back to generic
        openai_api_key = os.getenv("AZURE_OPENAI_KEY_41") or os.getenv("AZURE_OPENAI_KEY") or ""

        # ----------- Vectorizer עם api_key -----------
        vec_params = AzureOpenAIVectorizerParameters(
            resource_url    = azure_openai_endpoint,
            deployment_name = embedding_deployment,
            model_name      = embedding_model,
            api_key         = openai_api_key,
        )

        index_schema = SearchIndex(
            name   = name,
            fields = [
                SearchField(name="id", type="Edm.String", key=True, filterable=True, sortable=True, facetable=True),
                SearchableField(name="page_chunk", type="Edm.String", analyzer_name="standard.lucene"),
                SearchField(
                    name="page_embedding_text_3_large",
                    type="Collection(Edm.Single)",
                    stored=False,
                    vector_search_dimensions=VECTOR_DIM,
                    vector_search_profile_name="hnsw_text_3_large",
                ),
                SimpleField(name="page_number",  type="Edm.Int32",  filterable=True, sortable=True, facetable=True),
                SimpleField(name="source_file",  type="Edm.String", filterable=True, facetable=True),
                SimpleField(name="source",       type="Edm.String", filterable=True, facetable=True),
                SimpleField(name="url",          type="Edm.String"),
                SimpleField(name="doc_key",      type="Edm.String", filterable=True), # Added for proper document referencing
                # Enhanced metadata fields for Document Intelligence processing
                SimpleField(name="extraction_method", type="Edm.String", filterable=True, facetable=True),
                SimpleField(name="document_type", type="Edm.String", filterable=True, facetable=True),
                SimpleField(name="has_figures", type="Edm.Boolean", filterable=True, facetable=True),
                SimpleField(name="processing_timestamp", type="Edm.DateTimeOffset", filterable=True, sortable=True),
                # Multimodal fields for image processing
                SearchableField(name="content", type="Edm.String", analyzer_name="standard.lucene"),
                SearchField(name="contentVector",
                           type="Collection(Edm.Single)",
                           stored=False,
                           vector_search_dimensions=VECTOR_DIM,
                           vector_search_profile_name="hnsw_text_3_large"),
                SimpleField(name="imageCaptions", type="Edm.String", searchable=True, retrievable=True),
                SearchField(name="captionVector",
                           type="Collection(Edm.Single)",
                           stored=False,
                           vector_search_dimensions=VECTOR_DIM,
                           vector_search_profile_name="hnsw_text_3_large"),
                SimpleField(name="relatedImages", type="Collection(Edm.String)", filterable=True, retrievable=True),
                SimpleField(name="isMultimodal", type="Edm.Boolean", filterable=True, facetable=True),
                SimpleField(name="filename", type="Edm.String", filterable=True, facetable=True),
            ],
            vector_search = VectorSearch(
                profiles   = [ VectorSearchProfile(name="hnsw_text_3_large", algorithm_configuration_name="alg",
                                                   vectorizer_name="azure_open_ai_text_3_large") ],
                algorithms = [ HnswAlgorithmConfiguration(name="alg") ],
                vectorizers= [ AzureOpenAIVectorizer(vectorizer_name="azure_open_ai_text_3_large",
                                                     parameters=vec_params) ],           # ← משתמשים ב-vec_params
            ),
            semantic_search = SemanticSearch(
                default_configuration_name="semantic_config",
                configurations=[ SemanticConfiguration(
                    name="semantic_config",
                    prioritized_fields=SemanticPrioritizedFields(
                        content_fields=[ SemanticField(field_name="page_chunk") ]
                    ),
                )],
            ),
        )

        # מוחקים אינדקס קודם כדי לעדכן במקום
        if name in [idx.name for idx in index_client.list_indexes()]:
            index_client.delete_index(name)
        index_client.create_or_update_index(index_schema)

        # ----------- Knowledge-Agent עם api_key -------
        agent = KnowledgeAgent(
            name = f"{name}-agent",
            models = [
                KnowledgeAgentAzureOpenAIModel(
                    azure_open_ai_parameters = AzureOpenAIVectorizerParameters(
                        resource_url    = azure_openai_endpoint,
                        deployment_name = env("AZURE_OPENAI_DEPLOYMENT_41"),
                        model_name      = "gpt-4.1",
                        api_key         = openai_api_key,
                    )
                )
            ],
            target_indexes = [
                KnowledgeAgentTargetIndex(index_name=name, default_reranker_threshold=2.5)
            ],
        )
        index_client.create_or_update_agent(agent)
        return True

    except Exception as exc:
        st.error(f"Failed to create index '{name}': {exc}")
        return False


##############################################################################
# LLM prompts
##############################################################################

PLANNER_SYSTEM_PROMPT = textwrap.dedent(
    """
    You are a query‑planning assistant. Rewrite or split the **user question**
    Return **only** a JSON
    array of strings – no extra text.
    """
).strip()

ANSWER_SYSTEM_PROMPT = textwrap.dedent(
    """
    You are an AI assistant grounded in internal knowledge from Azure AI Search.
    • Use **only** the context passages below.  
    • When you quote, keep the exact citation label already inside the brackets – do **not** invent new labels.  
      Example: if the passage includes “[מב 50.02.pdf]” then cite exactly “[מב 50.02.pdf]”.  
    • If you lack information – say so honestly.  
    • Output in Markdown.
    """
).strip()

##############################################################################
# Agentic RAG core
##############################################################################


def plan_queries(question: str, client: AzureOpenAI, params: dict) -> List[str]:
    msgs = [
        {"role": "system", "content": PLANNER_SYSTEM_PROMPT},
        {"role": "user", "content": question},
    ]
    resp = client.chat.completions.create(messages=msgs, **params)
    txt = resp.choices[0].message.content.strip()
    try:
        return json.loads(txt)
    except json.JSONDecodeError:
        return [txt]


def retrieve(queries: List[str], client: SearchClient) -> List[dict]:
    docs: List[dict] = []
    for q in queries:
        hits = list(client.search(q, top=st.session_state.get("top_k", TOP_K_DEFAULT)))
        logging.warning("🔍 query='%s'  hits=%s", q, len(hits))
        for res in hits:
            docs.append(
                {
                    "id": len(docs) + 1,
                    "query": q,
                    "score": res["@search.score"],
                    "content": res.get("content", str(res))[:1000],
                }
            )
    return docs


def build_context(docs: List[dict]) -> str:
    """
    Build a concise context string by taking the first TOP‑K documents overall
    (no per‑query grouping) and truncating each passage to 600 characters.
    """
    chunk_size = 600
    top_k = (
        st.session_state.get("top_k", TOP_K_DEFAULT)
        if "top_k" in st.session_state
        else TOP_K_DEFAULT
    )

    return "\n\n".join(
        f"[doc{d['id']}] {d['content'][:chunk_size]}…" for d in docs[:top_k]
    )


def answer(question: str, ctx: str, client: AzureOpenAI, params: dict) -> tuple[str, int]:
    msgs = [
        {"role": "system", "content": ANSWER_SYSTEM_PROMPT},
        {"role": "system", "content": f"Context:\n\n{ctx}"},
        {"role": "user", "content": question},
    ]
    resp = client.chat.completions.create(messages=msgs, **params)
    answer_txt = resp.choices[0].message.content.strip()
    # `resp.usage` is a `CompletionUsage` object, not a dict
    tokens_used = getattr(resp, "usage", None)
    if tokens_used is not None and hasattr(tokens_used, "total_tokens"):
        tokens_used = tokens_used.total_tokens
    else:
        tokens_used = 0
    return answer_txt, tokens_used


##############################################################################
# CLI entry‑point
##############################################################################

##############################################################################
# Pipeline‑as‑a‑Tool helper: wraps KnowledgeAgentRetrievalClient.retrieve
# ---------------------------------------------------------------------------
def agentic_retrieval(agent_name: str, index_name: str, messages: list[dict]) -> str:
    """
    מבצע שליפה סוכנתית (Agentic Retrieval) באמצעות KnowledgeAgentRetrievalClient עבור agent נתון ו-index נתון.
    """
    # הגנה על פורמט ההודעות: המרה לסchema תקינה
    fixed_msgs = []
    for m in messages:
        if isinstance(m, dict) and "role" in m and "content" in m:
            fixed_msgs.append(m)
        elif isinstance(m, str):
            fixed_msgs.append({"role": "user", "content": m})
        else:
            raise ValueError(f"Unknown message format: {m}")

    ka_client = KnowledgeAgentRetrievalClient(
        endpoint=env("AZURE_SEARCH_ENDPOINT"),
        agent_name=agent_name,
        credential=_search_credential(),
    )
    ka_msgs = [
        KnowledgeAgentMessage(
            role=m["role"],
            content=[KnowledgeAgentMessageTextContent(text=m["content"])]
        )
        for m in fixed_msgs
    ]
    # ------------------ build retrieval request ---------------------------
    # If caller supplies an index_name, try to force the query there; otherwise
    # fall back to the agent’s default target index.
    target_params: list[KnowledgeAgentIndexParams] | None = None
    if index_name:
        target_params = [
            KnowledgeAgentIndexParams(
                index_name=index_name,
                reranker_threshold=2.5,
            )
        ]

    # Create a base request with only the most essential parameters
    req_params = {
        "messages": ka_msgs,
        # Only include target_index_params if we actually specified one
        "target_index_params": target_params,
        "request_limits": KnowledgeAgentRequestLimits(max_output_size=6000)
    }
    
    # Try to add optional parameters that might not be supported in all SDK versions
    try:
        # Create a test instance to check supported parameters
        test_req = KnowledgeAgentRetrievalRequest(messages=ka_msgs)
        
        # Check if citation_field_name is supported
        if hasattr(test_req, "citation_field_name"):
            req_params["citation_field_name"] = "source_file"
        
        # Add any other potentially unsupported parameters here
    except Exception:
        pass  # Silently continue with base parameters
    
    # Create the actual request
    req = KnowledgeAgentRetrievalRequest(**req_params)

    # ------------------ execute – retry without explicit index on mismatch -
    try:
        result = ka_client.knowledge_retrieval.retrieve(retrieval_request=req)
    except HttpResponseError as err:
        # If the agent is not configured for *index_name*, retry letting the
        # agent use its default target index.
        if (
            "target index name must match" in str(err).lower()
            and target_params is not None
        ):
            req.target_index_params = None  # remove the conflicting override
            result = ka_client.knowledge_retrieval.retrieve(retrieval_request=req)
        else:
            raise  # re‑raise unrelated errors

    # ----------------------------------------------------------------------
    # Build a rich JSON array with metadata so downstream agents can show
    # proper citations (url / source_file / page_number).
    # Each chunk inside `result.response` is a KnowledgeAgentMessage object
    # that holds one or more `content` items.
    # We flatten everything into a list like:
    #   [{"ref_id": 0, "content": "...", "url": "...", "source_file": "...", "page_number": 3}, …]
    # ----------------------------------------------------------------------
    chunks: list[dict] = []
    for msg in result.response:
        for c in getattr(msg, "content", []):
            chunk = {
                # ref_id might be absent – fall back to running index
                "ref_id": getattr(c, "ref_id", None) or len(chunks),
                "content": getattr(c, "text", ""),
                "url": getattr(c, "url", None),
                "source_file": getattr(c, "source_file", None),
                "page_number": getattr(c, "page_number", None),
                "score": getattr(c, "score", None),
                "doc_key": getattr(c, "doc_key", None),
            }
            # prune empty keys
            chunks.append({k: v for k, v in chunk.items() if v is not None})

    # Return the raw JSON string (no extra formatting)
    return json.dumps(chunks, ensure_ascii=False)

##############################################################################
# ZIP Function Folder Helper
##############################################################################
def _zip_function_folder(func_dir: Path, zip_path: Path) -> None:
    """
    Zip up the contents of *func_dir* (including all subfolders) into *zip_path*.
    Ensures all files are stored relative to func_dir (so host.json is at root).
    """
    import zipfile
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in func_dir.rglob("*"):
            if p.is_file():
                zf.write(p, p.relative_to(func_dir))
# ────────────────────────── Helper: zip Azure Function folder ──────────────
# def _zip_function_folder(func_dir: Path, zip_path: Path) -> None:
#     """Zip *func_dir* (כולל host.json וכו') כ-relative paths אל *zip_path*."""
#     with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
#         for itm in func_dir.rglob("*"):
#             if itm.is_file():
#                 zf.write(itm, itm.relative_to(func_dir))
# -----------------------------------------------------------------------------
# Processing Information Display Helper
# -----------------------------------------------------------------------------
def display_processing_info(file_name: str, file_ext: str, chunker_type: str = None, show_capabilities: bool = True):
    """
    Display processing information for a file to help users understand 
    what tools and methods are being used for extraction.
    """
    ext = file_ext.lower()
    
    # File type mapping
    file_type_map = {
        '.pdf': '📄 PDF Document',
        '.docx': '📝 Word Document', 
        '.pptx': '📊 PowerPoint Presentation',
        '.xlsx': '📈 Excel Spreadsheet',
        '.xls': '📈 Excel Spreadsheet',
        '.csv': '📈 CSV Data',
        '.png': '🖼️ PNG Image',
        '.jpg': '🖼️ JPEG Image',
        '.jpeg': '🖼️ JPEG Image',
        '.bmp': '🖼️ BMP Image',
        '.tiff': '🖼️ TIFF Image',
        '.txt': '📝 Text File',
        '.md': '📝 Markdown File',
        '.json': '🔧 JSON Data',
        '.html': '🌐 HTML Document',
        '.vtt': '🎬 Video Transcript'
    }
    
    # Processing method mapping
    processing_map = {
        '.pdf': ('🔍 Azure Document Intelligence', 'Advanced OCR, layout analysis, table extraction'),
        '.docx': ('🔍 Azure Document Intelligence', 'Layout analysis, text extraction, formatting preservation'),
        '.pptx': ('🔍 Azure Document Intelligence', 'Slide analysis, text extraction, layout understanding'),
        '.xlsx': ('🐼 Pandas Parser', 'Structured spreadsheet data extraction'),
        '.xls': ('🐼 Pandas Parser', 'Legacy Excel format processing'),
        '.csv': ('🐼 Pandas Parser', 'Comma-separated values processing'),
        '.png': ('🔍 Azure Document Intelligence', 'OCR text extraction from images'),
        '.jpg': ('🔍 Azure Document Intelligence', 'OCR text extraction from images'),
        '.jpeg': ('🔍 Azure Document Intelligence', 'OCR text extraction from images'),
        '.bmp': ('🔍 Azure Document Intelligence', 'OCR text extraction from images'),
        '.tiff': ('🔍 Azure Document Intelligence', 'OCR text extraction from images'),
        '.txt': ('📝 Simple Text Parser', 'Direct text content extraction'),
        '.md': ('📝 Markdown Parser', 'Markdown formatting with text extraction'),
        '.json': ('🔧 JSON Parser', 'Structured JSON data processing'),
        '.html': ('🔍 Azure Document Intelligence', 'HTML structure and content analysis'),
        '.vtt': ('🎬 Transcript Processor', 'Video subtitle and timing extraction')
    }
    
    file_type = file_type_map.get(ext, f'📄 {ext.upper()} File')
    method, capabilities = processing_map.get(ext, ('🔗 LangChain Chunker', 'General purpose text processing'))
    
    info_container = st.container()
    with info_container:
        col1, col2, col3 = st.columns([2, 3, 3])
        
        with col1:
            st.markdown(f"**File:** {file_type}")
            st.markdown(f"📋 `{file_name}`")
            
        with col2:
            st.markdown(f"**Processing Tool:** {method}")
            if show_capabilities:
                st.markdown(f"⚙️ {capabilities}")
                
        with col3:
            if ext in ['.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.html']:
                st.markdown("🎯 **Advanced Features:**")
                features = ["✅ Layout Analysis", "✅ Smart Text Extraction", "✅ OCR Processing"]
                if chunker_type == "MultimodalChunker":
                    features.extend(["✅ Figure Detection", "✅ AI Image Captions", "✅ Multimodal Processing"])
                for feature in features:
                    st.markdown(f"   {feature}")

# -----------------------------------------------------------------------------
# Streamlit UI wrapper (run with: streamlit run agentic-rag-demo.py)
# -----------------------------------------------------------------------------
def run_streamlit_ui() -> None:
    # Import required modules at function scope to avoid namespace conflicts
    import json as local_json
    from subprocess import check_output, CalledProcessError
    
    st.set_page_config(page_title="Agentic RAG Demo", page_icon="📚", layout="wide")

    # ── persistent session keys ───────────────────────────────────────────
    for k, default in {
        "selected_index": None,
        "available_indexes": [],
        "uploaded_files": [],
        "indexed_documents": {},
        "history": [],
        "agent_messages": [],
        "dbg_chunks": 0,
        "raw_index_json": "",  # last raw JSON from retrieval
        "orchestrator_targets": {},  # mapping: orchestrator agent → retrieval agent
    }.items():
        st.session_state.setdefault(k, default)

    st.title("📚 Agentic Retrieval‑Augmented Chat")
    st.markdown(
        """
        <style>
        html, body, .stApp { direction: ltr; text-align: left; }
        .ltr { direction: ltr; text-align: left; }
        </style>
        """,
        unsafe_allow_html=True,
    )

    # ── Sidebar – model & RAG knobs ───────────────────────────────────────
    with st.sidebar:
        st.header("⚙️ Model: GPT‑4.1")
        model_choice = "41"
        oai_client, chat_params = init_openai(model_choice)

        st.caption("Change `.env` to add more deployments")
        auth_mode = "Azure AD" if not os.getenv("AZURE_SEARCH_KEY") else "API Key"
        st.caption(f"🔑 Search auth: {auth_mode}")
        rbac_flag = _rbac_enabled(env("AZURE_SEARCH_ENDPOINT"))
        st.caption(f"🔒 RBAC: {'🟢 Enabled' if rbac_flag else '🔴 Disabled'}")
        if not rbac_flag:
            st.warning(
                "Turn on **Role‑based access control (Azure RBAC)** under "
                "*Search service → Networking → Authentication*."
            )

        st.subheader("🛠️ RAG Parameters")
        st.session_state.ctx_size = st.slider("Context chars per chunk", 300, 2000, 600, 50)
        st.session_state.top_k = st.slider("TOP‑K per query", 1, 200, 5, 1)
        st.session_state.rerank_thr = st.slider("Reranker threshold", 0.0, 4.0, 2.0, 0.1)
        st.session_state.max_output_size = st.slider("Knowledge‑agent maxOutputSize", 1000, 16000, 5000, 500)
        st.session_state.max_tokens = st.slider("Max completion tokens", 256, 32768, 32768, 256)

        chunks_placeholder = st.empty()
        chunks_placeholder.caption(f"Chunks sent to LLM: {st.session_state.get('dbg_chunks', 0)}")

        if st.button("🔄 Reload .env & restart"):
            _reload_env_and_restart()

    # ── Tabbed layout ─────────────────────────────────────────────────────
    # Initialize search client for index management
    _, root_index_client = init_search_client()
    
    tab_health, tab_create, tab_manage, tab_test, tab_cfg, tab_ai = st.tabs([
        "🩺 Health Check",
        "1️⃣ Create Index",
        "2️⃣ Manage Index",
        "3️⃣ Test Retrieval",
        "⚙️ Function Config",
        "🤖 AI Foundry Agent"
    ])

    # Health Check Tab
    with tab_health:
        # Initialize and render health check UI
        health_ui = HealthCheckUI()
        health_ui.render_health_check_tab()

    # Show warnings if health check not passed (optional, non-blocking)
    def health_block():
        health_ui = HealthCheckUI()
        health_ui.health_block()

    # ─────────────────── Tab 1 – Create Index ────────────────────────────
    with tab_create:
        health_block()
        st.header("🆕 Create a New Vector Index")
        new_index_name = st.text_input("New index name", placeholder="e.g. agentic‑vectors")
        if st.button("➕ Create new index") and new_index_name:
            if create_agentic_rag_index(root_index_client, new_index_name):
                st.success(f"Created index '{new_index_name}'")
                st.session_state.selected_index = new_index_name
                if new_index_name not in st.session_state.available_indexes:
                    st.session_state.available_indexes.append(new_index_name)

    # ─────────────────── Tab 2 – Manage Index ────────────────────────────
    with tab_manage:
        health_block()
        st.header("📂 Manage Existing Index")

        # refresh list each render
        st.session_state.available_indexes = [idx.name for idx in root_index_client.list_indexes()]

        existing = st.selectbox(
            "Existing indexes",
            options=[""] + st.session_state.available_indexes,
            index=0 if not st.session_state.selected_index else
                   st.session_state.available_indexes.index(st.session_state.selected_index)+1
                   if st.session_state.selected_index in st.session_state.available_indexes else 0,
            placeholder="Select index"
        )
        if existing:
            st.session_state.selected_index = existing
            st.success(f"Selected index: {existing}")

        # delete selected
        if st.session_state.selected_index:
            st.warning(f"Selected index: **{st.session_state.selected_index}**")
            if st.button("🗑️ Delete selected index"):
                try:
                    idx_name = st.session_state.selected_index
                    agent_name = f"{idx_name}-agent"
                    try:
                        root_index_client.delete_agent(agent_name)
                    except Exception:
                        pass
                    root_index_client.delete_index(idx_name)
                    st.session_state.available_indexes.remove(idx_name)
                    st.session_state.selected_index = None
                    st.success(f"Deleted index **{idx_name}** and its agent.")
                except Exception as ex:
                    st.error(f"Failed to delete index: {ex}")

        st.divider()
        st.subheader("📄 Upload PDFs into Selected Index")
        st.markdown(
            "פורמטים נתמכים בהעלאה ישירה: **PDF, DOCX, PPTX, XLSX/CSV, TXT, MD, JSON**  \n"
            "_קבצים אחרים יידחו אוטומטית או יועלו כ‑binary ללא חיפוש סמנטי._"
        )
        
        # Processing Information Section
        with st.expander("ℹ️ Document Processing Information", expanded=False):
            st.markdown("### 🔧 Processing Tools & Capabilities")
            
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("#### 🔍 Azure Document Intelligence")
                st.markdown("**Supported:** PDF, DOCX, PPTX, Images (PNG, JPG, BMP, TIFF)")
                st.markdown("**Capabilities:**")
                st.markdown("✅ Advanced OCR with high accuracy")
                st.markdown("✅ Layout and structure analysis")
                st.markdown("✅ Table extraction and formatting")
                st.markdown("✅ Figure and image detection")
                st.markdown("✅ Page-aware text chunking")
                st.markdown("✅ Multimodal processing (when enabled)")
                
                st.markdown("#### 🐼 Pandas Parser")
                st.markdown("**Supported:** CSV, XLS, XLSX")
                st.markdown("**Capabilities:**")
                st.markdown("✅ Structured data extraction")
                st.markdown("✅ Multiple sheet processing")
                st.markdown("✅ Data type preservation")
                
            with col2:
                st.markdown("#### 🔗 LangChain Chunker")
                st.markdown("**Supported:** General text files")
                st.markdown("**Capabilities:**")
                st.markdown("✅ Smart text chunking")
                st.markdown("✅ Overlap management")
                st.markdown("✅ Token-aware splitting")
                
                st.markdown("#### 📝 Simple Parser")
                st.markdown("**Supported:** TXT, MD, JSON")
                st.markdown("**Capabilities:**")
                st.markdown("✅ Direct text extraction")
                st.markdown("✅ Format preservation")
                st.markdown("✅ Fast processing")
            
            st.markdown("---")
            st.info("💡 **Tip:** Office documents (DOCX, PPTX) now automatically use Azure Document Intelligence for better structure preservation and metadata extraction!")
        
        if not st.session_state.selected_index:
            st.info("Select an index first.")
        else:
            uploaded = st.file_uploader(
                "בחר קבצים (PDF, DOCX, PPTX, XLSX/CSV, TXT, MD, JSON, RTX, XML)",
                type=["pdf", "docx", "pptx", "xlsx", "csv", "txt", "md", "json", "rtx", "xml"],
                accept_multiple_files=True
            )
            if uploaded and st.button("🚀 Ingest"):
                # Display processing overview
                st.markdown("### 🔄 Processing Overview")
                for pf in uploaded:
                    ext = os.path.splitext(pf.name)[-1].lower()
                    display_processing_info(pf.name, ext, show_capabilities=False)
                    st.markdown("---")
                
                with st.spinner("Embedding and uploading…"):
                    ###############################################
                    # Build buffered sender with error‑tracking
                    ###############################################
                    failed_ids: list[str] = []

                    def _on_error(action) -> None:
                        try:
                            # IndexAction object doesn't have .get() method, need to access attributes
                            if hasattr(action, 'id'):
                                failed_ids.append(action.id)
                            elif hasattr(action, 'document') and hasattr(action.document, 'get'):
                                failed_ids.append(action.document.get("id", "?"))
                            else:
                                failed_ids.append("?")
                        except Exception as exc:
                            logging.error("⚠️  on_error callback failed to record ID: %s", exc)
                            failed_ids.append("?")

                    sender = SearchIndexingBufferedSender(
                        endpoint=env("AZURE_SEARCH_ENDPOINT"),
                        index_name=st.session_state.selected_index,
                        credential=_search_credential(),
                        batch_size=100,
                        auto_flush_interval=5,
                        on_error=_on_error,
                    )

                    embed_deploy = os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT", "text-embedding-3-large")
                    total_pages = 0
                    processed_files = []
                    skipped_files = []
                    
                    for pf in uploaded:
                        ext = os.path.splitext(pf.name)[-1].lower()
                        docs = []

                        # --- Use DocumentChunker for ALL files including PDFs for multimodal support ---
                        error_message = None
                        
                        # DEBUG: Check file size at the very beginning
                        original_file_size = len(bytes(pf.getbuffer()))
                        logging.info(f"[Streamlit Upload][{pf.name}] ORIGINAL FILE SIZE: {original_file_size:,} bytes")
                        
                        # Also check the Streamlit file object properties
                        logging.info(f"[Streamlit Upload][{pf.name}] File object type: {type(pf)}")
                        logging.info(f"[Streamlit Upload][{pf.name}] File object size property: {getattr(pf, 'size', 'N/A')}")
                        
                        try:
                            docs = _chunk_to_docs(
                                pf.name,
                                bytes(pf.getbuffer()),
                                "",          # no public URL for local upload
                                oai_client,
                                embed_deploy,
                            )
                            
                            # Check if file was processed successfully
                            if not docs:
                                # Check if the original chunker returned any useful error information
                                multimodal_enabled = os.getenv("MULTIMODAL", "false").lower() in ["true", "1", "yes"] and ext in ('.pdf', '.png', '.jpeg', '.jpg', '.bmp', '.tiff', '.docx', '.pptx')
                                dc = DocumentChunker(multimodal=multimodal_enabled, openai_client=oai_client if multimodal_enabled else None)
                                data = {
                                    "fileName": pf.name,
                                    "documentBytes": base64.b64encode(bytes(pf.getbuffer())).decode("utf-8"),
                                    "documentUrl": "",
                                }
                                chunks, errors, warnings = dc.chunk_documents(data)
                                
                                if errors:
                                    error_message = f"Processing failed: {errors[0] if errors else 'Unknown error'}"
                                else:
                                    error_message = "No content could be extracted from this file"
                                    
                        except Exception as docerr:
                            error_message = str(docerr)
                            logging.error("DocumentChunker failed for %s: %s", pf.name, docerr)
                            
                            # Try fallback for PDFs only
                            if ext == ".pdf":
                                try:
                                    docs = pdf_to_documents(pf, oai_client, embed_deploy)
                                    error_message = None  # Clear error if fallback succeeded
                                    logging.info("Fallback to simple PDF processing for %s", pf.name)
                                except Exception as pdf_err:
                                    logging.error("PDF fallback also failed for %s: %s", pf.name, pdf_err)
                                    error_message = f"PDF processing failed: {str(docerr)[:200]}... (Fallback also failed: {str(pdf_err)[:100]}...)"
                            else:
                                # For non-PDF files, keep the original error
                                error_message = f"Failed to process {ext} file: {str(docerr)[:300]}..."
                        
                        # Handle errors - show in UI and track for summary
                        if error_message or not docs:
                            # Enhanced error message based on common issues
                            file_size = len(bytes(pf.getbuffer()))
                            enhanced_error = error_message or "Unknown processing error"
                            
                            # Provide specific guidance for common issues
                            guidance = ""
                            if file_size < 1000:
                                guidance = """
                                **This file is very small ({} bytes) which suggests it may be:**
                                - Corrupted or incomplete
                                - An empty file
                                - A file that failed to upload properly
                                
                                **Try:**
                                - Re-downloading the original file
                                - Checking if it opens properly in its native application
                                - Using a different version of the file
                                """.format(file_size)
                            elif "Document Intelligence" in enhanced_error and "UnsupportedContent" in enhanced_error:
                                guidance = """
                                **Document Intelligence couldn't process this file because:**
                                - The file may be corrupted or have invalid internal structure
                                - It might be password-protected
                                - The format may not be fully compatible
                                
                                **Try:**
                                - Opening and re-saving the file in its native application
                                - Converting to a different format (e.g., PDF → DOCX)
                                - Ensuring the file isn't password-protected
                                """
                            elif ext == ".pdf":
                                guidance = """
                                **PDF processing failed. Common causes:**
                                - Corrupted PDF file
                                - Password-protected PDF
                                - Non-standard PDF encoding
                                - Scanned PDF without OCR text layer
                                
                                **Try:**
                                - Re-saving the PDF from its source application
                                - Using a PDF repair tool
                                - Converting to Word format first
                                """
                            
                            st.error(f"""
                            **❌ Processing Failed: {pf.name}**
                            
                            {enhanced_error}
                            
                            **File details:**
                            - Size: {file_size} bytes  
                            - Type: {ext}
                            
                            {guidance}
                            
                            **What you can try:**
                            - Check if the file opens correctly in its native application
                            - Try re-saving or converting the file to a different format
                            - For PDFs: ensure they're not password-protected
                            - For images: ensure they're in a standard format
                            """)
                            
                            skipped_files.append({
                                "name": pf.name,
                                "size": len(bytes(pf.getbuffer())),
                                "reason": error_message or "Processing failed"
                            })
                            continue

                        if not docs:
                            skipped_files.append({
                                "name": pf.name,
                                "size": len(bytes(pf.getbuffer())),
                                "reason": "Corrupted or unsupported file"
                            })
                            continue
                            
                        # Show processing information to user
                        processing_info = []
                        for doc in docs[:1]:  # Check first document for processing info
                            method = doc.get("extraction_method", "unknown")
                            doc_type = doc.get("document_type", "Unknown")
                            has_figs = doc.get("has_figures", False)
                            
                            if method == "document_intelligence":
                                processing_info.append(f"📄 **{pf.name}** ({doc_type})")
                                processing_info.append("🔍 **Processing Tool:** Azure Document Intelligence")
                                processing_info.append("✨ **Capabilities:** Advanced layout analysis, OCR, table extraction")
                                if has_figs:
                                    processing_info.append("🖼️ **Figures:** Detected and processed with multimodal AI")
                            elif method == "simple_parser":
                                processing_info.append(f"📄 **{pf.name}** ({doc_type})")
                                processing_info.append("🔧 **Processing Tool:** Simple text parser")
                                processing_info.append("📝 **Capabilities:** Basic text extraction")
                            elif method == "pandas_parser":
                                processing_info.append(f"📊 **{pf.name}** ({doc_type})")
                                processing_info.append("🐼 **Processing Tool:** Pandas data parser")
                                processing_info.append("📈 **Capabilities:** Structured data extraction")
                            elif method == "langchain_chunker":
                                processing_info.append(f"📄 **{pf.name}** ({doc_type})")
                                processing_info.append("🔗 **Processing Tool:** LangChain document loader")
                                processing_info.append("⚡ **Capabilities:** Smart text chunking")
                                # Check if this was a fallback from multimodal/Document Intelligence
                                if ext in ('.pdf', '.png', '.jpeg', '.jpg', '.bmp', '.tiff', '.docx', '.pptx'):
                                    processing_info.append("⚠️ **Note:** Fell back to basic text extraction (Document Intelligence unavailable or file unsupported)")
                            else:
                                # Unknown method - show basic info
                                processing_info.append(f"📄 **{pf.name}** ({doc_type})")
                                processing_info.append(f"🔧 **Processing Tool:** {method}")
                        
                        # Add multimodal status info
                        multimodal_docs = [doc for doc in docs if doc.get("isMultimodal", False)]
                        if multimodal_docs:
                            processing_info.append(f"🎨 **Multimodal Content:** {len(multimodal_docs)} chunks contain images/figures")
                        elif ext in ('.pdf', '.png', '.jpeg', '.jpg', '.bmp', '.tiff') and os.getenv("MULTIMODAL", "false").lower() in ["true", "1", "yes"]:
                            processing_info.append("ℹ️ **Multimodal Status:** No images detected or multimodal processing failed")
                        
                        if processing_info:
                            with st.expander(f"ℹ️ Processing Details for {pf.name}", expanded=False):
                                for info in processing_info:
                                    st.markdown(info)
                                st.markdown(f"📊 **Chunks Created:** {len(docs)}")
                        
                        sender.upload_documents(documents=docs)
                        total_pages += len(docs)
                        processed_files.append({
                            "name": pf.name,
                            "chunks": len(docs),
                            "method": docs[0].get("extraction_method", "unknown") if docs else "unknown"
                        })

                    sender.close()

                    try:
                        search_client, _ = init_search_client(st.session_state.selected_index)
                        for _ in range(30):
                            if search_client.get_document_count() > 0:
                                break
                            time.sleep(1)
                    except Exception as probe_err:
                        logging.warning("Search probe failed: %s", probe_err)

                    success_pages = total_pages - len(failed_ids)
                    if failed_ids:
                        st.error(f"❌ {len(failed_ids)} pages failed to index – see logs for details.")
                    if success_pages:
                        st.success(f"✅ Indexed {success_pages} pages into **{st.session_state.selected_index}**.")
                    
                    # Show processing summary
                    if processed_files or skipped_files:
                        st.markdown("### 📊 Processing Summary")
                        
                        if processed_files:
                            st.markdown(f"**✅ Successfully Processed ({len(processed_files)} files):**")
                            for file_info in processed_files:
                                st.markdown(f"   • {file_info['name']} - {file_info['chunks']} chunks ({file_info['method']})")
                        
                        if skipped_files:
                            st.markdown(f"**⚠️ Skipped Files ({len(skipped_files)} files):**")
                            for file_info in skipped_files:
                                st.markdown(f"   • {file_info['name']} ({file_info['size']} bytes) - {file_info['reason']}")
                            st.info("💡 **Tip:** Skipped files are usually corrupted, too small, or in an unsupported format. Try re-saving or converting them.")

            # --- SharePoint Ingestion Button and Logic ---
            st.markdown("---")
            st.subheader("📁 Ingest files from SharePoint Folder (any type)")
            st.caption(
                "פורמטים מומלצים: **PDF, DOCX, PPTX, XLSX/CSV, TXT, MD, JSON**.  "
                "תוכל להזין ב‑תיבה למטה את סיומות קבצים מופרדות בפסיקים, או להשאיר ריק ל‑“הכול”."
            )

            # --- SharePoint config UI ---
            def _get_env_or_default(key, default=None):
                v = os.getenv(key)
                return v if v is not None and v != "" else default

            if "sp_site_domain" not in st.session_state:
                st.session_state.sp_site_domain = _get_env_or_default("SHAREPOINT_SITE_DOMAIN", "")
            if "sp_site_name" not in st.session_state:
                st.session_state.sp_site_name = _get_env_or_default("SHAREPOINT_SITE_NAME", "")
            if "sp_drive_name" not in st.session_state:
                st.session_state.sp_drive_name = _get_env_or_default("SHAREPOINT_DRIVE_NAME", "")
            if "sp_folder_path" not in st.session_state:
                st.session_state.sp_folder_path = _get_env_or_default("SHAREPOINT_SITE_FOLDER", "/")

            st.text("Current SharePoint settings:")
            st.session_state.sp_site_domain = st.text_input(
                "SharePoint Site Domain (e.g. mngenvmcap623661.sharepoint.com)",
                value=st.session_state.sp_site_domain or "",
                key="sp_site_domain_input"
            )
            st.session_state.sp_site_name = st.text_input(
                "SharePoint Site Name (blank or 'root' for root site)",
                value=st.session_state.sp_site_name or "",
                key="sp_site_name_input"
            )
            st.session_state.sp_drive_name = st.text_input(
                "SharePoint Drive Name (e.g. Documents)",
                value=st.session_state.sp_drive_name or "",
                key="sp_drive_name_input"
            )
            st.session_state.sp_folder_path = st.text_input(
                "SharePoint Folder Path (e.g. /ASKMANHAR)",
                value=st.session_state.sp_folder_path or "/",
                key="sp_folder_path_input"
            )

            file_type_input = st.text_input(
                "File types to ingest (comma-separated, e.g. pdf,docx,xlsx,txt). Leave blank for all:",
                value="pdf"
            )
            file_types = [ft.strip() for ft in file_type_input.split(",") if ft.strip()] if file_type_input else None
            if st.button("🔗 Ingest from SharePoint"):
                with st.spinner("Fetching and ingesting files from SharePoint…"):
                    # Initialize processing status display
                    status_container = st.empty()
                    
                    try:
                        from connectors.sharepoint.sharepoint_data_reader import SharePointDataReader
                        import io
                        sharepoint_reader = SharePointDataReader()
                        # Use user-specified values, fallback to env if blank
                        site_domain = st.session_state.sp_site_domain or _get_env_or_default("SHAREPOINT_SITE_DOMAIN", "")
                        site_name = st.session_state.sp_site_name or _get_env_or_default("SHAREPOINT_SITE_NAME", "")
                        folder_path = st.session_state.sp_folder_path or _get_env_or_default("SHAREPOINT_SITE_FOLDER", "/")
                        drive_name = st.session_state.sp_drive_name or _get_env_or_default("SHAREPOINT_DRIVE_NAME", "")
                        sp_files = sharepoint_reader.retrieve_sharepoint_files_content(
                            site_domain=site_domain,
                            site_name=site_name,
                            folder_path=folder_path,
                            file_formats=file_types,
                            drive_name=drive_name
                        )
                        if not sp_files:
                            st.warning("No files found in the SharePoint folder.")
                        else:
                            total_files = len(sp_files)
                            st.info(f"Found {total_files} file(s) in the SharePoint folder.")
                            progress_bar = st.progress(0, text="Starting ingestion...")
                            failed_ids = []
                            def _sp_on_error(action):
                                try:
                                    if hasattr(action, 'id'):
                                        failed_ids.append(action.id)
                                    elif hasattr(action, 'document') and hasattr(action.document, 'get'):
                                        failed_ids.append(action.document.get("id", "?"))
                                    else:
                                        failed_ids.append("?")
                                except Exception as exc:
                                    logging.error("⚠️  SharePoint on_error callback failed: %s", exc)
                                    failed_ids.append("?")
                            
                            sender = SearchIndexingBufferedSender(
                                endpoint=env("AZURE_SEARCH_ENDPOINT"),
                                index_name=st.session_state.selected_index,
                                credential=_search_credential(),
                                batch_size=100,
                                auto_flush_interval=5,
                                on_error=_sp_on_error,
                            )
                            embed_deploy = os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT", "text-embedding-3-large")
                            total_pages = 0
                            for idx, file in enumerate(sp_files):
                                file_bytes = file.get("content")
                                fname = file.get("name")
                                file_url = file.get("source") or file.get("webUrl")
                                st.info(f"Processing file {idx+1}/{total_files}: {fname}")
                                if not file_bytes or not fname:
                                    progress_bar.progress((idx + 1) / total_files, text=f"Skipped file {idx+1}/{total_files}")
                                    continue
                                
                                # If it's a PDF file
                                if fname.lower().endswith('.pdf'):
                                    pdf_file = io.BytesIO(file_bytes)
                                    pdf_file.name = fname
                                    docs = pdf_to_documents(pdf_file, oai_client, embed_deploy)
                                    # Patch each doc's url to SharePoint webUrl
                                    for d in docs:
                                        d["url"] = file_url
                                else:
                                    try:
                                        docs = _chunk_to_docs(
                                            fname,
                                            file_bytes,
                                            file_url,
                                            oai_client,
                                            embed_deploy,
                                        )
                                    except Exception as derr:
                                        logging.error("Chunker failed for %s: %s", fname, derr)
                                        docs = []
                                sender.upload_documents(documents=docs)
                                total_pages += len(docs)
                                progress_bar.progress((idx + 1) / total_files, text=f"Processed {idx+1}/{total_files} files")
                            sender.close()
                            try:
                                search_client, _ = init_search_client(st.session_state.selected_index)
                                for _ in range(30):
                                    if search_client.get_document_count() > 0:
                                        break
                                    time.sleep(1)
                            except Exception as probe_err:
                                logging.warning("Search probe failed: %s", probe_err)
                            success_pages = total_pages - len(failed_ids)
                            if failed_ids:
                                st.error(f"❌ {len(failed_ids)} pages failed to index – see logs for details.")
                            if success_pages:
                                st.success(f"✅ Indexed {success_pages} pages from SharePoint into **{st.session_state.selected_index}**.")
                    except Exception as ex:
                        st.error(f"SharePoint ingestion failed: {ex}")

    # ─────────────────── Tab 3 – Test Retrieval ──────────────────────────
    with tab_test:
        render_test_retrieval_tab(
            tab_test, 
            health_block, 
            st.session_state, 
            init_agent_client, 
            init_search_client, 
            env, 
            _search_credential
        )

    # ─────────────────── Tab 5 – Function Config ──────────────────────────
    with tab_cfg:
        st.header("⚙️ Configure Azure Function")

        # Load .env once for this tab
        env_vars = load_env_vars()

        # Runtime parameters – choose index from dropdown
        st.markdown("##### Runtime parameters")

        # Populate from earlier tabs (Manage/Create) – falls back to manual
        index_options = st.session_state.get("available_indexes", [])
        if index_options:
            # Pre‑select value from .env if present
            try:
                preselect = index_options.index(env_vars.get("INDEX_NAME", index_options[0]))
            except ValueError:
                preselect = 0
            idx_selected = st.selectbox("INDEX_NAME", index_options, index=preselect)
        else:
            st.warning("No index list detected – enter manually.")
            idx_selected = st.text_input("INDEX_NAME", env_vars.get("INDEX_NAME", ""))

        # Update env_vars with the chosen/typed value
        env_vars["INDEX_NAME"] = idx_selected.strip()
        env_vars["AGENT_NAME"] = f"{idx_selected.strip()}-agent" if idx_selected else ""

        # Display the derived AGENT_NAME (read‑only)
        st.text_input("AGENT_NAME", env_vars["AGENT_NAME"], disabled=True)

        # Try to pre‑fill subscription from az cli
        cli_sub = get_azure_subscription()
        sub_id = st.text_input("Subscription ID", cli_sub)

        # List Function Apps in this subscription
        func_choices, func_map = list_function_apps(sub_id)
        
        if not func_choices and sub_id:
            st.warning("⚠️ Could not list Function Apps automatically; fill manually.")

        func_sel_lbl = st.selectbox(
            "Choose Function App",
            ["-- manual input --"] + func_choices,
            index=0
        )
        st.session_state["func_map"] = func_map
        st.session_state["func_choices"] = func_choices
        
        if func_sel_lbl != "-- manual input --":
            app, rg = func_map[func_sel_lbl]
        else:
            rg = st.text_input("Resource Group", os.getenv("AZURE_RG", ""))
            app = st.text_input("Function App name", os.getenv("AZURE_FUNCTION_APP", ""))
        
        # Normalise variable names (func_name / func_rg) and keep old aliases
        func_name = app
        func_rg = rg

        if not all((sub_id, rg, app)):
            st.info("Fill subscription / RG / Function-App and click 🔄 Load settings.")
        else:
            if "func_raw" not in st.session_state:
                st.session_state.func_raw = {}
            if "func_df" not in st.session_state:
                st.session_state.func_df = pd.DataFrame(columns=["key", "value"])

            if st.button("🔄 Load settings"):
                success, df, raw, error_msg = load_function_settings(rg, app, sub_id, env_vars)
                if success:
                    st.session_state.func_raw = raw
                    st.session_state.func_df = df
                    st.success(f"Loaded & merged {len(df)} setting(s).")
                else:
                    st.error(f"Failed to load: {error_msg}")

        # Show editable table on every render once loaded
        if st.session_state.get("func_df") is not None and not st.session_state.func_df.empty:
            st.markdown("#### Function App Settings")
            st.session_state.func_df = _st_data_editor(
                st.session_state.func_df,
                num_rows="dynamic",
                use_container_width=True,
                key="func_editor",
            )

            # Push edited settings back to the Function App
            st.divider()
            if st.button("💾 Push settings to Function"):
                success, message = push_function_settings(
                    func_rg, 
                    func_name, 
                    sub_id, 
                    st.session_state.func_df,
                    st.session_state.func_raw
                )
                if success:
                    st.success(f"✅ {message} on **{func_name}**")
                else:
                    st.error(f"Failed to update Function settings:\n{message}")

        # Deploy local ./function code to this Function App
        st.divider()
        if st.button("🚀 Deploy local code to Function"):
            with st.spinner("⏳ Zipping and deploying, please wait…"):
                success, message, stdout = deploy_function_code(func_rg, func_name, sub_id)
                if success:
                    st.success(f"✅ {message}")
                    if stdout:
                        st.text(stdout)
                else:
                    st.error(message)

    # ─────────────────── Tab 6 – AI Foundry Agent ────────────────────────
    with tab_ai:
        health_block()
        st.header("🤖 Create AI Foundry Agent")

        func_map     = st.session_state.get("func_map", {})
        func_choices = st.session_state.get("func_choices", [])

        if not func_choices:
            st.info("Go to **Function Config** tab first and load settings.")
            st.stop()

        func_sel = st.selectbox("Function App to invoke", func_choices, index=0)
        func_name, func_rg = func_map[func_sel]
        base_url = f"https://{func_name}.azurewebsites.net/api"

        # Detect Foundry projects the CLI user can access
        cli_cred = AzureCliCredential()
        logged_in, _ = check_azure_cli_login()
        if not logged_in:
            st.error("🔑 Run `az login` before using this feature.")
            st.stop()

        projects = get_ai_foundry_projects(cli_cred)
        
        # If no projects found, show helpful message
        if not projects:
            st.warning("No AI Foundry projects found via Azure CLI.")
            
            # Check for PROJECT_ENDPOINT in .env
            project_endpoint_env = os.getenv("PROJECT_ENDPOINT", "").strip()
            if project_endpoint_env:
                projects = [{
                    "name": project_endpoint_env.split('/')[-1][:30] or "env-project",
                    "location": "env",
                    "endpoint": project_endpoint_env,
                    "resource_group": "env",
                    "hub_name": "env",
                }]
                st.success(f"Using PROJECT_ENDPOINT from .env: {project_endpoint_env}")
            else:
                # Allow manual entry
                st.info(
                    "You can either:\n"
                    "1. Create a project in Azure AI Studio\n"
                    "2. Set the PROJECT_ENDPOINT environment variable in your .env file\n"
                    "3. Make sure you have access to at least one AI Foundry project"
                )
                
                manual_endpoint = st.text_input(
                    "Or enter Project Endpoint manually:",
                    placeholder="https://my-project.api.region.ai.azure.com/"
                )
                if manual_endpoint:
                    projects = [{
                        "name": "manual-project",
                        "location": "manual",
                        "endpoint": manual_endpoint,
                        "resource_group": "manual",
                        "hub_name": "manual",
                    }]
                else:
                    st.stop()

        proj_labels = [f"{p['name']} – {p['location']}" for p in projects]
        sel = st.selectbox("Choose Foundry project", proj_labels, index=0)
        project_endpoint = projects[proj_labels.index(sel)]['endpoint']
        st.caption(f"🔗 Endpoint: {project_endpoint}")

        agent_name = st.text_input("Agent name", placeholder="function‑assistant")
        if st.button("🚀 Create Agent") and agent_name:
            # Use the refactored function instead of inline code
            success, message, agent = create_ai_foundry_agent(
                project_endpoint=project_endpoint,
                agent_name=agent_name,
                base_url=base_url,
                function_key=FUNCTION_KEY
            )
            
            if success:
                st.success(f"✅ Agent **{agent.name}** created (ID: {agent.id})")
            else:
                st.error("Failed to create agent via SDK:")
                st.error(message)


##############################################################################
# Main entry point
##############################################################################

def main():
    """Main CLI entry point."""
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "question",
        nargs="?",
        help="Your question (e.g. 'What is Azure AI Search?')"
    )
    parser.add_argument(
        "--model",
        choices=["o3", "4o", "41"],
        default="41",
        help="Model variant: o3 | 4o | 41 (default: 41)"
    )
    args = parser.parse_args()

    # If no question provided or running under Streamlit, show UI
    if not args.question or _st_in_runtime():
        run_streamlit_ui()
        return

    # CLI mode - process the question
    oai_client, chat_params = init_openai(args.model)
    search_client, _ = init_search_client(os.getenv("INDEX_NAME", "agentic-vectors"))

    # Plan queries
    queries = plan_queries(args.question, oai_client, chat_params)
    print(f"📋 Planned queries: {queries}")

    # Retrieve documents
    docs = retrieve(queries, search_client)
    print(f"📚 Retrieved {len(docs)} documents")

    # Build context and answer
    ctx = build_context(docs)
    final_answer, tokens = answer(args.question, ctx, oai_client, chat_params)
    
    print("\n" + "="*80)
    print(f"💬 Answer ({tokens} tokens):\n")
    print(final_answer)
    print("="*80 + "\n")


if __name__ == "__main__":
    main()